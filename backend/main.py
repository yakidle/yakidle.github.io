import os
import uuid
from fastapi import FastAPI, File, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, FileResponse
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
# 你需要实现或者替换以下两行
# from openai_api import call_chat_model
# from utils import extract_key_images

app = FastAPI()
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])

UPLOAD_DIR = "uploads"
PPT_DIR = "ppt_files"
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(PPT_DIR, exist_ok=True)

def extract_text_from_pdf(pdf_path):
    text = ""
    doc = fitz.open(pdf_path)
    for page in doc:
        text += page.get_text()
    return text

# 假定大模型API如下，建议你换成自己的API/KEY
def call_chat_model(prompt):
    # 你可以接第三方API (例如OpenAI、Qwen等)：此处为伪代码
    # response = openai.ChatCompletion.create(...)
    # return response['choices'][0]['message']['content']
    return "【这里是大模型返回的论文解析内容】"

def extract_key_images(pdf_path, limit=1):
    # 返回第1张图片作为示意
    doc = fitz.open(pdf_path)
    images = []
    for page in doc:
        for img in page.get_images(full=True):
            xref = img[0]
            image_bytes = doc.extract_image(xref)["image"]
            fname = f"{uuid.uuid4().hex[:8]}.png"
            fout = os.path.join(PPT_DIR, fname)
            with open(fout, "wb") as f:
                f.write(image_bytes)
            images.append(fout)
            if len(images) >= limit:
                break
        if len(images) >= limit:
            break
    return images

def generate_ppt(summary:str, images:list, filename:str):
    prs = Presentation()
    # 封面
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "智能论文讲解"
    slide.placeholders[1].text = "自动生成PPT摘要"
    # 核心讲解
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "讲解摘要"
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(7), Inches(3))
    tf = txBox.text_frame
    tf.text = summary[:600]
    # 核心图片
    for i, img in enumerate(images):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"论文关键流程图{(i+1)}"
        slide.shapes.add_picture(img, Inches(1.5), Inches(1.8), width=Inches(5))
    fout = os.path.join(PPT_DIR, filename)
    prs.save(fout)
    return fout

@app.post("/upload")
async def upload(pdf: UploadFile = File(...)):
    temp = os.path.join(UPLOAD_DIR, f"{uuid.uuid4().hex[:8]}.pdf")
    with open(temp, "wb") as f:
        f.write(await pdf.read())
    text = extract_text_from_pdf(temp)
    # 1. 用大模型判断类型, 这里只做方法类
    type_prompt = f"请判断以下论文属于哪一类型（如综述/方法/项目汇报/理论推导/应用等）：\n{text[:1600]}"
    paper_type = call_chat_model(type_prompt)
    # 2. 如果是方法类，用方法类prompt生成结构化讲解
    main_prompt = (
        f"你是一位AI论文讲解专家，请针对以下论文内容，判断其模型架构与工作流程，用简明语言讲解主要创新点、待解决的问题和局限性，"
        "最后分3~4步摘要整份论文的工作流程，每步50字左右。内容：\n" + text[:3000]
    )
    summary = call_chat_model(main_prompt)
    # 3. 抽取一张图片
    key_imgs = extract_key_images(temp, limit=1)
    # 4. 生成PPT
    ppt_fn = f"{uuid.uuid4().hex[:8]}.pptx"
    ppt_path = generate_ppt(summary, key_imgs, ppt_fn)
    ppt_url = f"/ppt/{ppt_fn}"
    return {"success": True, "paper_type": paper_type, "summary": summary, "ppt_url": ppt_url}

@app.get("/ppt/{ppt_file}")
def get_ppt(ppt_file:str):
    ppt_path = os.path.join(PPT_DIR, ppt_file)
    if not os.path.isfile(ppt_path):
        return JSONResponse({"success":False,"message":"文件不存在"},status_code=404)
    return FileResponse(ppt_path, filename=ppt_file, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
